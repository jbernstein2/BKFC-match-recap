import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from branding import COLORS


# ───────────────────────────────
# Helpers
# ───────────────────────────────

def rgb(hex_color):
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=COLORS["BLACK"]):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.text = str(text)

    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)

    return box


# ───────────────────────────────
# TITLE SLIDE (MATCHES YOUR NODE VERSION)
# ───────────────────────────────

def add_title_slide(prs, data):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["BLACK"])

    # gold top bar
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS["GOLD"])

    add_text(
        slide,
        "BROOKLYN FC",
        0.5, 0.8, 12, 1,
        size=54,
        bold=True,
        color=COLORS["GOLD"]
    )

    add_text(
        slide,
        "MATCH REPORT",
        0.5, 1.8, 12, 0.5,
        size=16,
        bold=True,
        color=COLORS["WHITE"]
    )

    if data["score"]:
        add_text(
            slide,
            data["score"],
            5, 2.5, 3, 1,
            size=32,
            bold=True,
            color=COLORS["GOLD"]
        )

    add_text(
        slide,
        f"BKFC vs {data['opponent_name']}",
        0.5, 3.2, 12, 0.6,
        size=24,
        bold=True,
        color=COLORS["WHITE"]
    )

    add_text(
        slide,
        data["match_date"],
        0.5, 3.9, 12, 0.4,
        size=14,
        color=COLORS["GOLD"]
    )

    add_text(
        slide,
        data["competition"],
        0.5, 4.3, 12, 0.4,
        size=11,
        color=COLORS["SILVER"]
    )


# ───────────────────────────────
# SUMMARY SLIDE (clean table style)
# ───────────────────────────────

def add_summary_slide(prs, data):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, "FFFFFF")

    add_text(slide, "MATCH SUMMARY", 0.5, 0.3, 10, 0.5, size=20, bold=True)

    metrics = [
        ("Goals", 6),
        ("xG", 7),
        ("Shots", 8),
        ("Possession %", 14)
    ]

    y = 1.2

    for label, col in metrics:

        bkfc_val = round(float(data["match_bkfc"][col]), 2)
        opp_val = round(float(data["match_opp"][col]), 2)

        add_text(slide, label, 0.5, y, 3, 0.4, size=14, bold=True)
        add_text(slide, bkfc_val, 4, y, 2, 0.4, size=14)
        add_text(slide, opp_val, 6, y, 2, 0.4, size=14)

        y += 0.6


# ───────────────────────────────
# INSIGHTS SLIDE
# ───────────────────────────────

def add_insights_slide(prs, insights):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["BLACK"])

    add_text(slide, "KEY INSIGHTS", 0.5, 0.5, 10, 0.5,
             size=26, bold=True, color=COLORS["GOLD"])

    y = 1.5

    for i in insights:

        add_text(
            slide,
            "• " + i,
            0.7, y, 10, 0.5,
            size=14,
            color=COLORS["WHITE"]
        )

        y += 0.5


# ───────────────────────────────
# STAT SLIDE (clean 2-column layout)
# ───────────────────────────────

def add_stat_slide(prs, label, match_vals, season_vals, opp_name):

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, "F4F6F9")

    add_text(slide, label.upper(), 0.5, 0.3, 10, 0.5,
             size=18, bold=True, color=COLORS["BLACK"])

    # Match
    add_text(slide, "THIS MATCH", 1, 1, 3, 0.4, size=12, bold=True)
    add_text(slide, "BKFC", 1, 1.5, 3, 0.4, size=14)
    add_text(slide, opp_name, 1, 2, 3, 0.4, size=14)

    add_text(slide, match_vals["bkfc"], 2.5, 1.5, 2, 0.4, size=14)
    add_text(slide, match_vals["opp"], 2.5, 2, 2, 0.4, size=14)

    # Season
    add_text(slide, "SEASON AVG", 6, 1, 3, 0.4, size=12, bold=True)
    add_text(slide, "BKFC", 6, 1.5, 3, 0.4, size=14)
    add_text(slide, "League", 6, 2, 3, 0.4, size=14)
    add_text(slide, opp_name, 6, 2.5, 3, 0.4, size=14)

    add_text(slide, season_vals["bkfcAvg"], 8, 1.5, 2, 0.4, size=14)
    add_text(slide, season_vals["allOppAvg"], 8, 2, 2, 0.4, size=14)
    add_text(slide, season_vals["thisOpp"], 8, 2.5, 2, 0.4, size=14)


# ───────────────────────────────
# MAIN GENERATOR
# ───────────────────────────────

def generate_report(data, output_file, stats, insights):

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, data)
    add_summary_slide(prs, data)
    add_insights_slide(prs, insights)

    for stat in stats:

        add_stat_slide(
            prs,
            stat["label"],
            stat["match"],
            stat["season"],
            data["opponent_name"]
        )

    prs.save(output_file)

    return output_file
